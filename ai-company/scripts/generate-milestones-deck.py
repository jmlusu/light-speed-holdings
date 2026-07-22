#!/usr/bin/env python3
"""Generate AI Company Builder Milestones Deck — full 15-slide version."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

C = {
    "dk": RGBColor(0x1A, 0x1A, 0x2E),
    "tl": RGBColor(0x16, 0x21, 0x3E),
    "ac": RGBColor(0x0F, 0x34, 0x60),
    "hl": RGBColor(0xE9, 0x45, 0x60),
    "wh": RGBColor(0xFF, 0xFF, 0xFF),
    "lg": RGBColor(0xF8, 0xF9, 0xFA),
    "mg": RGBColor(0x6C, 0x75, 0x7D),
    "dg": RGBColor(0x34, 0x3A, 0x40),
    "ok": RGBColor(0x28, 0xA7, 0x45),
}


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def dark_bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = C["dk"]


def white_bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = C["wh"]


def header_bar(slide, text):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C["dk"]
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["wh"]


def bullet_box(slide, x, y, w, h, title, items, title_size=16, item_size=13):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = C["dg"]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(item_size)
        p.font.color.rgb = C["dg"]
        p.space_before = Pt(4)


def text_line(slide, x, y, w, text, size=14, color=None, bold=False, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color or C["dg"]
    p.font.bold = bold
    p.font.italic = italic


def metric_card(slide, x, y, value, label, color):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.1), Inches(1.1))
    card.fill.solid()
    card.fill.fore_color.rgb = C["wh"]
    card.line.color.rgb = C["mg"]
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.07), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(1.8), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    tb2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.6), Inches(1.8), Inches(0.4))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = C["mg"]


def milestone_slide(prs, num, title, date, items):
    s = blank(prs)
    white_bg(s)
    header_bar(s, f"Milestone {num}: {title}")
    text_line(s, 0.5, 1.0, 9, date, size=12, color=C["mg"], italic=True)
    bullet_box(s, 0.5, 1.4, 9, 3.8, "", items, item_size=13)
    return s


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # --- Slide 1: Title ---
    s = blank(prs)
    dark_bg(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "AI Company Builder"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C["wh"]
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Major Milestones"
    p2.font.size = Pt(28)
    p2.font.color.rgb = C["wh"]
    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.4))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "v0.3.0  |  1205 Tests  |  26 CLI Commands  |  53 Agent Roles"
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = C["wh"]
    tb4 = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.3))
    p4 = tb4.text_frame.paragraphs[0]
    p4.text = "July 2026"
    p4.font.size = Pt(12)
    p4.font.color.rgb = C["mg"]

    # --- Slide 2: Executive Summary ---
    s = blank(prs)
    white_bg(s)
    header_bar(s, "Executive Summary")
    metric_card(s, 0.5, 1.1, "1205", "Tests Passing", C["ok"])
    metric_card(s, 2.8, 1.1, "26", "CLI Commands", C["ac"])
    metric_card(s, 5.1, 1.1, "53", "Agent Roles", C["hl"])
    metric_card(s, 7.4, 1.1, "0", "Lint Errors", C["ok"])
    bullet_box(s, 0.5, 2.5, 9, 2.8, "Project Overview", [
        "- Python CLI tool for creating and orchestrating AI agent hierarchies",
        "- 3 sprints completed, 8 major milestones achieved",
        "- 6 domain engines: Executor, Decision, Workflow, Memory, Graph, Dashboard",
        "- Production-ready with comprehensive testing (1205 tests) and documentation (30+ docs)",
    ])

    # --- Slide 3: Timeline ---
    s = blank(prs)
    white_bg(s)
    header_bar(s, "Project Timeline")
    milestones = [
        ("Foundation", "Pre-July"),
        ("Sprint 1", "Jul 20"),
        ("Sprint 2", "Jul 21"),
        ("Org Expansion", "Jul 21"),
        ("Sprint 3", "Jul 22"),
        ("Dashboard", "Jul 17+"),
        ("CLI (26 cmds)", "Ongoing"),
        ("Documentation", "Jul 17+"),
    ]
    for i, (label, date) in enumerate(milestones):
        x = 0.3 + i * 1.2
        dot = s.shapes.add_shape(1, Inches(x + 0.4), Inches(2.0), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C["hl"] if i in (2, 4) else C["ac"]
        dot.line.fill.background()
        text_line(s, x, 2.3, 1.1, label, size=9, bold=True)
        text_line(s, x, 2.65, 1.1, date, size=8, color=C["mg"])
    # line
    line = s.shapes.add_shape(1, Inches(0.3), Inches(2.08), Inches(9.4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C["mg"]
    line.line.fill.background()
    bullet_box(s, 0.5, 3.2, 9, 2.0, "Sprint Velocity", [
        "- Sprint 1: Code hardening + audit trail",
        "- Sprint 2: 13 items, 41 hrs -- 1093 tests passing",
        "- Sprint 3: 8 items, 22 hrs -- 1205 tests passing, v0.3.0 tagged",
    ])

    # --- Slide 4-11: Milestones ---
    milestone_slide(prs, 1, "Foundation", "Pre-July 2026", [
        "- Core Python package structure with Typer CLI framework",
        "- 17+ Pydantic domain models (Company, Executive, Department, Agent, Task, etc.)",
        "- 4-module registry system: loader, parser, resolver, validator (19 YAML configs)",
        "- 12 Jinja2 templates for agent generation (base, executive, department, specialist, etc.)",
        "- Generator with template selection by agent type",
        "- Organization: 7 executive roles, 5 core departments",
    ])

    milestone_slide(prs, 2, "Code Hardening & Audit Trail", "Sprint 1 -- 2026-07-20", [
        "- Audit trail system: AuditEvent, AuditWriter (JSONL append), AuditReader (query/filter)",
        "- Executor integration: tool calls, task lifecycle, HITL decisions logged",
        "- Code hardening across all modules (Track B)",
        "- Audit trail completion (Track C)",
        "- All Sprint 1 items verified complete",
    ])

    milestone_slide(prs, 3, "Core Engines & Integration", "Sprint 2 -- 2026-07-21", [
        "- Executor pipeline: Inbox polling -> AgentLoop (ReAct) -> LLM -> Tools -> Audit",
        "- AgentLoop: Multi-turn LLM<->tool, cost tracking, HITL gates, budget enforcement",
        "- DecisionEngine: Approval matrix, risk assessment, decision tree navigation",
        "- WorkflowEngine: 9 workflows, step tracking, SLA monitoring, task conversion",
        "- MemoryEngine: 6 types (episodic, semantic, procedural, relational, temporal, aggregate)",
        "- GraphEngine: 4 types (org_chart, decision, workflow, knowledge) with BFS pathfinding",
        "- Circuit Breaker, Cost Tracker (JSONL), Dead-Letter Queue, Postmortem system",
        "- 13 items completed, 41 hrs effort, 1093 tests passing",
    ])

    milestone_slide(prs, 4, "Organization Expansion", "2026-07-21", [
        "- 53 new roles added across all departments",
        "- Phase 1 (Immediate): VP Eng, Data Engineer, AI Safety Lead, Program Manager, etc.",
        "- Phase 2 (Weeks 1-8): Red Team, MLOps, Platform Eng, Product Marketing, etc.",
        "- Phase 3 (Weeks 9-16): Frontend/API Architect, Observability, Privacy Officer, etc.",
        "- Phase 4 (Weeks 17-24): Technical Writer, L&D Lead, IR Lead, BI Engineer, etc.",
        "- Structural improvements: CTO span reduced 9+ -> 1, Product dept (8 roles), Strategy dept",
        "- AI Safety hierarchy: AI Safety Lead -> Red Team, Constitutional AI, Ethics",
    ])

    milestone_slide(prs, 5, "Gap Closure & Testing", "Sprint 3 -- 2026-07-22", [
        "- Org chart test rewrite: 832 lines -> 56 tests, all passing",
        "- DataTransformer.registry_to_enhanced() frozen model bug fix",
        "- E2E pipeline test covering full executor workflow",
        "- WebSocket integration tests (30 tests) for real-time broadcasts",
        "- Governance CLI: 7 commands (report, audit-trail, risk-summary, retention, compliance, owners, policies)",
        "- Memory CLI enhancement: stats, search, recall commands",
        "- Dashboard API tests (9 tests) for CEO dashboard, KPIs, departments",
        "- 8 items completed, 22 hrs effort, 1205 tests, v0.3.0 tagged",
    ])

    milestone_slide(prs, 6, "Dashboard & Monitoring", "2026-07-17+", [
        "- FastAPI REST API with CORS and API key authentication",
        "- WebSocket broadcast for real-time task/KPI/alert updates",
        "- 7-department KPI system: Engineering, HR, Finance, Legal, Marketing, Sales, Customer Success",
        "- 28 KPI definitions with automated collectors",
        "- Analytics engine: History tracking, trend analysis, alert rules, summary rollups",
        "- CEO Dashboard: company health, agent performance, cost tracking, task pipeline",
        "- Department dashboards with drill-down capability",
    ])

    milestone_slide(prs, 7, "CLI & Developer Experience", "Ongoing", [
        "- 26 CLI commands registered across all modules",
        "- Commands: company, decision, graph, workflows, memory, agents, board, departments,",
        "  executives, specialists, orchestrator, models, dashboard, executor, doctor,",
        "  marketing, sales, customer-success, legal, hr, generate, status, sop, raci, governance",
        "- System diagnostics (doctor command): Python version, dependencies, agent files,",
        "  inbox health, memory engine, disk space, cost tracker, LLM providers",
        "- Pre-commit hooks: ruff, mypy, bandit, trailing-whitespace, end-of-file-fixer, check-yaml",
    ])

    milestone_slide(prs, 8, "Documentation & Governance", "2026-07-17+", [
        "- 30+ documentation files covering architecture, development, ECL",
        "- Architecture docs, developer guide, ECL (change lifecycle)",
        "- SOPs: Incident response, deployment, HR onboarding, budget approval",
        "- RACI matrices: Hiring, escalation, deployment workflows",
        "- Risk register: 14 items with mitigations and owners",
        "- Board governance charter, meeting cadence, voting rules",
        "- Model routing policy: Provider catalog, tiers, routing rules, cost control",
        "- Company constitution: Principles and decision order",
    ])

    # --- Slide 12: Quality Metrics ---
    s = blank(prs)
    white_bg(s)
    header_bar(s, "Quality Metrics")
    rows = [
        ("Metric", "Value", "Status"),
        ("Pytest", "1205 tests passing", "GREEN"),
        ("Ruff (lint)", "0 errors", "GREEN"),
        ("Mypy (type check)", "0 errors (164 files)", "GREEN"),
        ("Bandit (security)", "No high-severity issues", "GREEN"),
        ("Version", "v0.3.0 tagged 2026-07-22", "RELEASED"),
        ("Pre-commit hooks", "6 hooks active", "ENABLED"),
    ]
    for i, (m, v, st) in enumerate(rows):
        y = 1.1 + i * 0.5
        bg = C["dk"] if i == 0 else (C["lg"] if i % 2 == 0 else C["wh"])
        fc = C["wh"] if i == 0 else C["dg"]
        rect = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(9), Inches(0.45))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.fill.background()
        text_line(s, 0.6, y + 0.07, 3, m, size=12, color=fc, bold=(i == 0))
        text_line(s, 3.6, y + 0.07, 4, v, size=12, color=fc, bold=(i == 0))
        sc = C["ok"] if st in ("GREEN", "RELEASED", "ENABLED") else fc
        text_line(s, 7.6, y + 0.07, 2, st, size=12, color=sc, bold=(i == 0))

    # --- Slide 13: Architecture ---
    s = blank(prs)
    white_bg(s)
    header_bar(s, "Architecture Overview")
    arch = [
        "cli/              26 Typer commands",
        "executor/         Task execution pipeline (AgentLoop, HITL, DLQ)",
        "llm/              Multi-provider LLM abstraction, circuit breaker, cost tracker",
        "orchestrator/     Message bus, scheduler, escalation, approval, tier rules",
        "models/           17+ Pydantic domain models",
        "registry/         4-module config registry (loader/parser/resolver/validator)",
        "builder/          BootstrapEngine -- full company generation",
        "decision/         DecisionEngine -- approvals, risk, decision trees",
        "workflow/         WorkflowEngine -- 9 workflows, SLA monitoring",
        "memory/           MemoryEngine -- 6 memory types, persistence, executor integration",
        "graph/            GraphEngine -- 4 graph types, BFS pathfinding",
        "audit/            Audit trail -- events, writer, reader, executor integration",
        "dashboard/        FastAPI REST API, WebSocket, KPI collectors (7 depts), analytics",
        "org_chart/        OrganizationChart -- O(1) lookup, pathfinding, subtree extraction",
    ]
    bullet_box(s, 0.5, 1.0, 9, 4.2, "", arch, item_size=11)

    # --- Slide 14: Remaining Work ---
    s = blank(prs)
    white_bg(s)
    header_bar(s, "Remaining Work -- Sprint 4")
    items = [
        ("Structured logging (GAP-018)", "High", "Correlation IDs, structured log format"),
        ("Scheduled cycle daemon mode", "High", "Persistent scheduler for autonomous cycles"),
        ("Agent spec validation CLI", "Medium", "Validate agent specs against schema"),
        ("CLI type hints/docstrings", "Medium", "Full type coverage for all 26 commands"),
        ("OAuth2/key rotation", "Medium", "Secure credential management"),
        ("Memory encryption", "Medium", "At-rest encryption for memory store"),
        ("Token counting integration", "Low", "Real-time token usage tracking"),
    ]
    for i, (item, pri, desc) in enumerate(items):
        y = 1.1 + i * 0.58
        bg = C["lg"] if i % 2 == 0 else C["wh"]
        rect = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(9), Inches(0.52))
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg
        rect.line.fill.background()
        pc = C["hl"] if pri == "High" else (RGBColor(0xFF, 0xC1, 0x07) if pri == "Medium" else C["mg"])
        text_line(s, 0.6, y + 0.05, 3, item, size=11, bold=True)
        text_line(s, 3.6, y + 0.05, 1, pri, size=10, color=pc, bold=True)
        text_line(s, 4.7, y + 0.05, 4.5, desc, size=10, color=C["mg"])

    # --- Slide 15: Next Steps ---
    s = blank(prs)
    dark_bg(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Next Steps"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C["wh"]
    items = [
        "1. Begin Sprint 4 with structured logging and daemon mode",
        "2. Full integration test of executor pipeline with real LLM calls",
        "3. Deploy to staging environment for end-to-end validation",
        "4. Target v0.4.0 release after Sprint 4 completion",
        "5. Plan Sprint 5 for external deployment and scaling",
    ]
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(3))
    tf = tb2.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = C["wh"]
        p.space_before = Pt(12)

    # Footer
    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.4))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "AI Company Builder v0.3.0  |  July 2026"
    p3.font.size = Pt(12)
    p3.font.color.rgb = C["mg"]

    # Save
    out = os.path.join(os.path.dirname(__file__), "..", "docs", "milestones-deck.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"Presentation created: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
