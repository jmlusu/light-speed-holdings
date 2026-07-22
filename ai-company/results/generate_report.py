#!/usr/bin/env python3
"""Generate AI Company Builder Progress Report PowerPoint deck.

Usage:
    python results/generate_report.py

Creates results/progress_report.pptx with a dark blue executive theme.
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ─────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0B, 0x19, 0x29)  # slide background
BG_CARD     = RGBColor(0x12, 0x2B, 0x45)  # card / section bg
BG_CARD_ALT = RGBColor(0x0E, 0x23, 0x3B)  # alternate card
ACCENT      = RGBColor(0x00, 0xD4, 0xAA)  # teal accent
ACCENT_DIM  = RGBColor(0x00, 0x9E, 0x7E)  # dimmer teal
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB8, 0xC9, 0xE0)
MID_GRAY    = RGBColor(0x7A, 0x8F, 0xA8)
DARK_TEXT    = RGBColor(0x0B, 0x19, 0x29)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, x, y, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, x, y, w, h, color=ACCENT):
    """Thin accent line."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_textbox(slide, x, y, w, h, text, font_size=14, color=WHITE,
                bold=False, align=PP_ALIGN.LEFT, font_name="Calibri",
                valign=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    tf.auto_size = None
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    return txBox


def add_bullet_list(slide, x, y, w, h, items, font_size=13,
                    color=LIGHT_GRAY, bullet_color=ACCENT, spacing=6):
    """Add a bulleted list. Items are strings."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
        # Use a run with colored bullet char + text
        run_bullet = p.add_run()
        run_bullet.text = "\u25CF  "  # filled circle bullet
        run_bullet.font.size = Pt(font_size - 1)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
    return txBox


def add_stat_card(slide, x, y, w, h, number, label, num_color=ACCENT):
    """Large stat number + small label card."""
    card = add_shape_rect(slide, x, y, w, h, BG_CARD)
    # Number
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4),
                Inches(0.6), number, font_size=36, color=num_color,
                bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
    # Label
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4),
                Inches(0.4), label, font_size=11, color=MID_GRAY,
                align=PP_ALIGN.CENTER, font_name="Calibri")
    return card


# ── Slide builders ───────────────────────────────────────────────────────────

def build_slide_01_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Decorative top accent bar
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))

    # Large decorative rectangle behind title area
    add_shape_rect(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.2), BG_CARD)

    # Title
    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(7.6), Inches(1.0),
                "AI Company Builder", font_size=42, color=WHITE, bold=True,
                align=PP_ALIGN.LEFT, font_name="Georgia")

    # Accent bar under title
    add_accent_bar(slide, Inches(1.2), Inches(2.35), Inches(2.5), Inches(0.05))

    # Subtitle line 1
    add_textbox(slide, Inches(1.2), Inches(2.6), Inches(7.6), Inches(0.5),
                "Progress Report", font_size=28, color=ACCENT,
                bold=False, align=PP_ALIGN.LEFT, font_name="Georgia")

    # Subtitle line 2
    add_textbox(slide, Inches(1.2), Inches(3.3), Inches(7.6), Inches(0.5),
                "Sprints 3-4 & Phase 3C Complete  |  July 21, 2026",
                font_size=15, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Bottom accent bar
    add_accent_bar(slide, Inches(0), Inches(5.565), SLIDE_W, Inches(0.06))


def build_slide_02_executive_summary(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Executive Summary", font_size=26, color=WHITE, bold=True,
                font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(1.8), Inches(0.04))

    # 5 stat cards across top
    card_w = Inches(1.68)
    card_h = Inches(1.1)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    y_cards = Inches(1.15)

    stats = [
        ("4/4", "Architecture\nGaps Closed"),
        ("30", "New QA\nTests"),
        ("962", "Tests\nPassing"),
        ("0", "Mypy\nErrors"),
        ("120", "Agents\nFrozen"),
    ]
    for i, (num, label) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        add_stat_card(slide, x, y_cards, card_w, card_h, num, label)

    # Key points below cards
    points = [
        "All critical architecture gaps (GAP-005, 014, 015, 020) closed",
        "30 new QA tests across approval escalation, scheduler, and LLM retry",
        "Semantic search fixed \u2014 VectorStore now properly initialized with EmbeddingEngine",
        "Codebase healthy: 962 tests passing, 0 mypy errors, 0 ruff violations",
        "120-agent registry frozen per CEO directive",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(2.5), Inches(8.8), Inches(2.8),
                    points, font_size=13, spacing=5)


def build_slide_03_architecture_fixes(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Architecture Fixes (GAPs)", font_size=26, color=WHITE,
                bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(2.0), Inches(0.04))

    # 2x2 grid of GAP cards
    gaps = [
        ("GAP-005", "Memory Consolidation Wired",
         "ConsolidationScheduler with tick + time triggers integrated into executor loop"),
        ("GAP-014", "BriefingGenerator DI",
         "BriefingGenerator now accepts dependency-injected MessageBus for testability"),
        ("GAP-015", "LLM Retry Cycling Fixed",
         "Flattened stream loop with provider rotation prevents infinite retry storms"),
        ("GAP-020", "E2E Pipeline Tests",
         "10 comprehensive end-to-end tests covering the full agent lifecycle"),
    ]

    card_w = Inches(4.3)
    card_h = Inches(1.8)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    start_x = Inches(0.5)
    start_y = Inches(1.1)

    for i, (gap_id, title, desc) in enumerate(gaps):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Card background
        add_shape_rect(slide, x, y, card_w, card_h, BG_CARD)
        # Left accent bar
        add_accent_bar(slide, x, y, Inches(0.06), card_h, ACCENT)

        # GAP ID badge
        add_textbox(slide, x + Inches(0.25), y + Inches(0.15), Inches(1.2),
                    Inches(0.3), gap_id, font_size=12, color=ACCENT,
                    bold=True, font_name="Consolas")

        # Title
        add_textbox(slide, x + Inches(0.25), y + Inches(0.45), card_w - Inches(0.5),
                    Inches(0.35), title, font_size=15, color=WHITE, bold=True)

        # Description
        add_textbox(slide, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5),
                    Inches(0.85), desc, font_size=11, color=MID_GRAY)


def build_slide_04_phase_3a(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Phase 3A \u2014 Frontend / Dashboard", font_size=26,
                color=WHITE, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(2.5), Inches(0.04))

    # Left column: bullet points
    points = [
        "OpenAPI / Swagger confirmed at /docs endpoint",
        "WebSocket broadcast wired for real-time task updates",
        "Rate limit headers added (X-RateLimit-Limit, X-RateLimit-Remaining)",
        "Dashboard fully operational with API key auth",
        "CORS configuration for cross-origin access",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(1.2), Inches(5.0), Inches(3.5),
                    points, font_size=14, spacing=8)

    # Right column: status card
    card_x = Inches(6.0)
    card_y = Inches(1.2)
    card_w = Inches(3.5)
    card_h = Inches(3.6)
    add_shape_rect(slide, card_x, card_y, card_w, card_h, BG_CARD)
    add_accent_bar(slide, card_x, card_y, card_w, Inches(0.05), ACCENT)

    add_textbox(slide, card_x + Inches(0.3), card_y + Inches(0.2),
                card_w - Inches(0.6), Inches(0.4),
                "STATUS", font_size=14, color=ACCENT, bold=True,
                font_name="Consolas")

    status_items = [
        ("OpenAPI", "READY"),
        ("WebSocket", "WIRED"),
        ("Rate Limits", "ACTIVE"),
        ("Auth", "ENABLED"),
        ("Dashboard", "OPERATIONAL"),
    ]
    for j, (item, status) in enumerate(status_items):
        iy = card_y + Inches(0.7) + j * Inches(0.55)
        add_textbox(slide, card_x + Inches(0.3), iy,
                    Inches(1.8), Inches(0.3),
                    item, font_size=12, color=LIGHT_GRAY)
        status_color = ACCENT if status in ("READY", "OPERATIONAL", "ACTIVE", "ENABLED") else MID_GRAY
        add_textbox(slide, card_x + Inches(2.2), iy,
                    Inches(1.0), Inches(0.3),
                    status, font_size=12, color=status_color, bold=True,
                    align=PP_ALIGN.RIGHT, font_name="Consolas")


def build_slide_05_phase_3b(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Phase 3B \u2014 ML / Memory", font_size=26,
                color=WHITE, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(2.0), Inches(0.04))

    # Before / After comparison
    col_w = Inches(4.3)
    gap = Inches(0.4)

    # BEFORE card
    bx = Inches(0.5)
    by = Inches(1.15)
    add_shape_rect(slide, bx, by, col_w, Inches(3.8), BG_CARD)
    add_accent_bar(slide, bx, by, col_w, Inches(0.05), RGBColor(0xFF, 0x6B, 0x6B))

    add_textbox(slide, bx + Inches(0.25), by + Inches(0.15), Inches(2), Inches(0.35),
                "BEFORE (Broken)", font_size=14,
                color=RGBColor(0xFF, 0x6B, 0x6B), bold=True)

    before_points = [
        "init_memory() created VectorStore without EmbeddingEngine",
        "Semantic search was dead code \u2014 substring fallback only",
        "No embedding model loaded, no vector index built",
        "Memory consolidation not wired into executor loop",
    ]
    add_bullet_list(slide, bx + Inches(0.25), by + Inches(0.55),
                    col_w - Inches(0.5), Inches(3.0),
                    before_points, font_size=12,
                    color=MID_GRAY, bullet_color=RGBColor(0xFF, 0x6B, 0x6B),
                    spacing=6)

    # AFTER card
    ax = bx + col_w + gap
    add_shape_rect(slide, ax, by, col_w, Inches(3.8), BG_CARD)
    add_accent_bar(slide, ax, by, col_w, Inches(0.05), ACCENT)

    add_textbox(slide, ax + Inches(0.25), by + Inches(0.15), Inches(2), Inches(0.35),
                "AFTER (Fixed)", font_size=14,
                color=ACCENT, bold=True)

    after_points = [
        "EmbeddingEngine(all-MiniLM-L6-v2) properly initialized",
        "Full vector index with index_all() for semantic search",
        "Memory consolidation integrated into executor tick loop",
        "6 memory types operational with persistence",
    ]
    add_bullet_list(slide, ax + Inches(0.25), by + Inches(0.55),
                    col_w - Inches(0.5), Inches(3.0),
                    after_points, font_size=12, spacing=6)


def build_slide_06_phase_3c(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Phase 3C \u2014 QA / Testing", font_size=26,
                color=WHITE, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(2.0), Inches(0.04))

    # 3 cards for the 3 test suites
    cards = [
        ("16 Tests", "Approval Escalation",
         ["HITL parking / unparking", "Approval rules engine",
          "Timeout handling", "Rejection flow", "Pre-approved actions"]),
        ("7 Tests", "Scheduler Verification",
         ["Init and config", "Tick increment logic",
          "Consolidation stats", "Time-based triggers"]),
        ("7 Tests", "LLM Retry Logic",
         ["Bad JSON retry loop", "Provider exhaustion",
          "GAP-015 stream fix", "Circuit breaker integration"]),
    ]

    card_w = Inches(2.9)
    card_h = Inches(3.6)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.1)

    for i, (count, title, items) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        # Card bg
        add_shape_rect(slide, x, start_y, card_w, card_h, BG_CARD)
        # Top accent
        add_accent_bar(slide, x, start_y, card_w, Inches(0.05), ACCENT)

        # Count
        add_textbox(slide, x + Inches(0.25), start_y + Inches(0.15),
                    card_w - Inches(0.5), Inches(0.4),
                    count, font_size=24, color=ACCENT, bold=True,
                    font_name="Georgia", align=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, x + Inches(0.25), start_y + Inches(0.55),
                    card_w - Inches(0.5), Inches(0.35),
                    title, font_size=14, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

        # Divider
        add_accent_bar(slide, x + Inches(0.8), start_y + Inches(0.95),
                       card_w - Inches(1.6), Inches(0.02), BG_CARD_ALT)

        # Items
        add_bullet_list(slide, x + Inches(0.25), start_y + Inches(1.1),
                        card_w - Inches(0.5), Inches(2.3),
                        items, font_size=11, spacing=5, color=LIGHT_GRAY)


def build_slide_07_quality_metrics(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Quality Metrics", font_size=26,
                color=WHITE, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(1.8), Inches(0.04))

    # 5 large stat callouts
    metrics = [
        ("962", "Tests Total", "(81 new this session)"),
        ("0", "Mypy Errors", "(was 3)"),
        ("0", "Ruff Violations", "(clean)"),
        ("120", "Registry Agents", "(frozen)"),
        ("4/4", "GAPs Closed", "(all resolved)"),
    ]

    card_w = Inches(1.68)
    card_h = Inches(1.8)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    y_cards = Inches(1.15)

    for i, (num, label, sublabel) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        add_shape_rect(slide, x, y_cards, card_w, card_h, BG_CARD)
        # Accent top
        add_accent_bar(slide, x, y_cards, card_w, Inches(0.04), ACCENT)

        # Big number
        add_textbox(slide, x + Inches(0.1), y_cards + Inches(0.2),
                    card_w - Inches(0.2), Inches(0.7),
                    num, font_size=40, color=ACCENT, bold=True,
                    align=PP_ALIGN.CENTER, font_name="Georgia")

        # Label
        add_textbox(slide, x + Inches(0.1), y_cards + Inches(0.9),
                    card_w - Inches(0.2), Inches(0.35),
                    label, font_size=13, color=WHITE,
                    align=PP_ALIGN.CENTER, bold=True)

        # Sublabel
        add_textbox(slide, x + Inches(0.1), y_cards + Inches(1.25),
                    card_w - Inches(0.2), Inches(0.3),
                    sublabel, font_size=10, color=MID_GRAY,
                    align=PP_ALIGN.CENTER)

    # Bottom summary line
    add_shape_rect(slide, Inches(0.5), Inches(3.3), Inches(9.0), Inches(1.8), BG_CARD)
    add_accent_bar(slide, Inches(0.5), Inches(3.3), Inches(0.06), Inches(1.8), ACCENT)

    summary_points = [
        "All 962 tests passing with zero failures across unit, integration, and E2E suites",
        "Type safety improved: mypy errors reduced from 3 to 0 in this session",
        "Code style clean: zero ruff lint violations, pre-commit hooks passing",
        "120-agent registry frozen per CEO directive \u2014 no additions or removals",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(3.4), Inches(8.5), Inches(1.6),
                    summary_points, font_size=12, spacing=4)


def build_slide_08_files(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Files Modified / Created", font_size=26,
                color=WHITE, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(2.2), Inches(0.04))

    # Two columns: NEW (left) and FIXED (right)
    col_w = Inches(4.3)
    gap = Inches(0.4)

    # NEW column
    nx = Inches(0.5)
    ny = Inches(1.1)
    add_shape_rect(slide, nx, ny, col_w, Inches(3.8), BG_CARD)
    add_accent_bar(slide, nx, ny, col_w, Inches(0.05), ACCENT)

    add_textbox(slide, nx + Inches(0.25), ny + Inches(0.15), Inches(2), Inches(0.35),
                "NEW FILES", font_size=14, color=ACCENT, bold=True,
                font_name="Consolas")

    new_files = [
        "test_approval_escalation.py  (16 tests)",
        "test_scheduler_verification.py  (7 tests)",
        "test_llm_retry_verification.py  (7 tests)",
        "test_full_pipeline.py  (10 E2E tests)",
    ]
    add_bullet_list(slide, nx + Inches(0.25), ny + Inches(0.55),
                    col_w - Inches(0.5), Inches(3.0),
                    new_files, font_size=12, spacing=7,
                    bullet_color=ACCENT)

    # FIXED column
    fx = nx + col_w + gap
    add_shape_rect(slide, fx, ny, col_w, Inches(3.8), BG_CARD)
    add_accent_bar(slide, fx, ny, col_w, Inches(0.05),
                   RGBColor(0xFF, 0xB8, 0x47))

    add_textbox(slide, fx + Inches(0.25), ny + Inches(0.15), Inches(2), Inches(0.35),
                "FIXED FILES", font_size=14,
                color=RGBColor(0xFF, 0xB8, 0x47), bold=True,
                font_name="Consolas")

    fixed_files = [
        "memory/consolidation.py",
        "llm/client.py",
        "orchestrator/briefing.py",
        "memory/integration.py",
        "dashboard/app.py",
        "cli/memory.py",
    ]
    add_bullet_list(slide, fx + Inches(0.25), ny + Inches(0.55),
                    col_w - Inches(0.5), Inches(3.0),
                    fixed_files, font_size=12, spacing=5,
                    bullet_color=RGBColor(0xFF, 0xB8, 0x47))


def build_slide_09_unfinished(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), BG_CARD)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.6),
                "Unfinished Work & Next Steps", font_size=26,
                color=WHITE, bold=True, font_name="Georgia")
    add_accent_bar(slide, Inches(0.6), Inches(0.75), Inches(2.5), Inches(0.04))

    # Warning callout box
    warn_x = Inches(0.5)
    warn_y = Inches(1.1)
    warn_w = Inches(9.0)
    warn_h = Inches(1.1)
    add_shape_rect(slide, warn_x, warn_y, warn_w, warn_h, BG_CARD)
    add_accent_bar(slide, warn_x, warn_y, Inches(0.06), warn_h,
                   RGBColor(0xFF, 0xB8, 0x47))

    add_textbox(slide, warn_x + Inches(0.25), warn_y + Inches(0.1),
                warn_w - Inches(0.5), Inches(0.3),
                "CEO DIRECTIVE", font_size=13,
                color=RGBColor(0xFF, 0xB8, 0x47), bold=True,
                font_name="Consolas")

    add_textbox(slide, warn_x + Inches(0.25), warn_y + Inches(0.45),
                warn_w - Inches(0.5), Inches(0.55),
                "Registry rationalization: DO NOT reduce \u2014 keep at 120 agents. "
                "No new features in pipeline until CEO direction.",
                font_size=13, color=LIGHT_GRAY)

    # Next steps
    points = [
        "Registry rationalization: CEO says DO NOT reduce \u2014 keep at 120 agents",
        "Next sprint priorities need CEO input",
        "No new features in pipeline until CEO direction",
        "Awaiting guidance on next phase of development",
        "Sprint 3 (Autonomous Coordination) and Sprint 4 (Quality) awaiting green light",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(2.5), Inches(8.8), Inches(2.8),
                    points, font_size=14, spacing=7)


def build_slide_10_closing(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Top accent bar
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))

    # Central card
    add_shape_rect(slide, Inches(1.5), Inches(1.0), Inches(7.0), Inches(3.6), BG_CARD)

    # Accent bar on card
    add_accent_bar(slide, Inches(1.5), Inches(1.0), Inches(7.0), Inches(0.05), ACCENT)

    # Main message
    add_textbox(slide, Inches(2.0), Inches(1.5), Inches(6.0), Inches(0.8),
                "Ready to Proceed", font_size=38, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, font_name="Georgia")

    add_accent_bar(slide, Inches(4.2), Inches(2.25), Inches(1.6), Inches(0.04))

    add_textbox(slide, Inches(2.0), Inches(2.5), Inches(6.0), Inches(0.5),
                "on CEO Direction", font_size=24, color=ACCENT,
                align=PP_ALIGN.CENTER, font_name="Georgia")

    # Contact info
    add_textbox(slide, Inches(2.0), Inches(3.4), Inches(6.0), Inches(0.4),
                "AI Company Builder Team", font_size=14, color=MID_GRAY,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.0), Inches(3.8), Inches(6.0), Inches(0.4),
                "962 Tests  |  120 Agents  |  0 Errors", font_size=12,
                color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Bottom accent bar
    add_accent_bar(slide, Inches(0), Inches(5.565), SLIDE_W, Inches(0.06))


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 10 slides
    build_slide_01_title(prs)
    build_slide_02_executive_summary(prs)
    build_slide_03_architecture_fixes(prs)
    build_slide_04_phase_3a(prs)
    build_slide_05_phase_3b(prs)
    build_slide_06_phase_3c(prs)
    build_slide_07_quality_metrics(prs)
    build_slide_08_files(prs)
    build_slide_09_unfinished(prs)
    build_slide_10_closing(prs)

    # Ensure output directory exists
    output_dir = Path(__file__).parent
    output_path = output_dir / "progress_report.pptx"
    output_dir.mkdir(parents=True, exist_ok=True)

    prs.save(str(output_path))
    print(f"Presentation saved to: {output_path.resolve()}")
    print(f"  {len(prs.slides)} slides generated")
    print("  Theme: Dark blue executive")
    print("  Layout: 16:9")


if __name__ == "__main__":
    main()
